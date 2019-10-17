from pptx import Presentation
prs = Presentation("H:/learning_notes/make_ppt_by_python/example.pptx")
for s in prs.slides:
    try:
        print(s.shapes.title.text)
    except:
        pass
    if s.shapes.placeholders:
        for p in s.shapes.placeholders:
            try:
                print(p.text)
            except:
                pass
