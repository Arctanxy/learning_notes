# https://blog.csdn.net/u011932355/article/details/73287248
from pptx import Presentation
from pptx.util import Inches,Cm,Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION

import requests as req
import re
from lxml import html
from urllib.request import urlretrieve
headers = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 6.1; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/67.0.3396.87 Safari/537.36',
}
r = req.get("https://www.jianshu.com/p/2ff1ead2249a",headers=headers)
tree = html.fromstring(r.text)
content = tree.xpath('//h1|//h2|//h4|//p|//img/@data-original-src')

prs = Presentation("H:/learning_notes/make_ppt_by_python/example.pptx")

# 创建封面ppt
slide = prs.slides.add_slide(prs.slide_layouts[0]) # slide_layouts[0] 为带有正标题和副标题的幻灯片
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = content[0].text # content中第一个元素是h1标题
subtitle.text = 'dalalaa'

slide1 = prs.slides.add_slide(prs.slide_layouts[1]) # slide_layouts[1] 为带有标题和文本框的幻灯片
title1 = slide1.shapes.title
textbox = slide1.placeholders[1]
title1.text = "简介"
textbox.text = content[1].text + content[2].text
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(5)
for i,c in enumerate(content[3:-12]):
    try:
        # 如果是二级标题，则创建新的标题+文本框页
        if '</h2>' in html.tostring(c).decode('utf-8'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = c.text 
        # 如果是四级标题，则创建新的标题+文本框页
        elif '</h4>' in html.tostring(c).decode('utf-8'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = c.text 
        # 如果是正文，则添加到文本框中
        elif '</p>' in html.tostring(c).decode('utf-8'):
            try:
                textbox = slide.placeholders[1]
            except:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                textbox = slide.placeholders[1]
            textbox.text += c.text
    except Exception as e:
        print(c)
        if 'upload' in c:
            filename = "H:/learning_notes/make_ppt_by_python/" + str(i) + ".jpg"
            urlretrieve("http:" + c,filename=filename)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            pic = slide.shapes.add_picture(filename, left, top,width,height)
        else:
            pass

prs.save("H:/learning_notes/make_ppt_by_python/example.pptx")


'''
        # 如果是正文，且上一个元素是h2标题，则另起一个空白页，在文本框中填写正文
        elif '</p>' in html.tostring(c).decode('utf-8') and '</h2>' in html.tostring(content[i-1]).decode('utf-8'):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(left,top,width,height)
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = c.text
            p.font.size = Pt(40)
            
        # 如果是正文，且上一个元素是h4标题，则将正文加入到文本框中
        elif '</p>' in html.tostring(c).decode('utf-8') and '</h4>' in html.tostring(content[i-1]).decode('utf-8'):
            textbox = slide.placeholders[1]
            textbox.text = c.text
        # 如果上一个元素也是正文，则将正文添加到前段文字后面
        elif '</p>' in html.tostring(c).decode('utf-8') and '</p>' in html.tostring(content[i-1]).decode('utf-8'):
            p = tf.add_paragraph()
            p.text = c.text
            p.font.size = Pt(40)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(left,top,width,height)
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = c.text
            p.font.size = Pt(40)
'''

